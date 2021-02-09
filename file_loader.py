from func.pdf_extractor import extract_pdf2str, convert_html2text
import struct

import docx2txt
from pptx import Presentation

accepted_filetypes = ['pdf', 'html'] + ['doc', 'docx'] + ['ppt', 'pptx'] + ['txt', 'py', 'cs', 'c', 'cpp', 'h', 'ipynb', 'hpp']

binary_filetypes = ['doc', 'docx', 'ppt', 'pptx', 'pdf']

def loadfile(fp, filetype, txt_split_size=1000, verbose=0, handle_unknown='skip'):

    allowed_cases = ['skip', 'upload_text_only', 'upload_all', 'raise']
    assert handle_unknown in allowed_cases, f"Value for handle_unknown was wrong: {handle_unknown}, but expected one of {allowed_cases}"  

    if filetype == 'pdf':
        pages_raw = extract_pdf2str(fp, verbose=verbose)
        page_nos, page_txts = convert_html2text(pages_raw, is_converted_pdf=True)
        data = dict(zip(page_nos, page_txts))
    elif filetype == 'html':
        pages_raw = {0: fp.load()}
        data = dict(zip(convert_html2text(pages_raw, is_converted_pdf=False)))
    elif filetype in ['doc', 'docx']:
        data = {0: docx2txt.process(fp)}
    elif filetype in ['ppt', 'pptx']:
        data = {}
        for i, slide in enumerate(Presentation(fp).slides):
            data[i] = '\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
    elif filetype in ['txt', 'py', 'cs', 'c', 'cpp', 'h', 'ipynb', 'hpp']:
        txt = fp.read()
        data = {cnt:txt[i:i+txt_split_size] for cnt, i in enumerate(range(0, len(txt), txt_split_size))}
    else:
        if handle_unknown == 'skip':
            pass
        elif handle_unknown == 'upload_all':
            if 'b' in fp.mode:
                fileContent = fp.read()
                fileContentStr = struct.unpack("i" * ((len(fileContent) -24) // 4), fileContent[20:-4])
                data = {0: fileContentStr}
            else:
                data = {0: fp.read()}
        elif handle_unknown == 'upload_text_only': 
            data = {0: fp.read()}
        elif handle_unknown == 'raise':
            raise ValueError("unknown type discovered with extension " + filetype)

    return data

